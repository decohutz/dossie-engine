"""
PowerPoint exporter for dossier data.
Generates a professional executive presentation with charts, tables, and KPI cards.
Palette: Midnight Executive (navy + ice blue + white).
"""
from __future__ import annotations
import io
import re
import unicodedata
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from ..models.dossier import Dossier

# ═══════════════════════════════════════════════════════════════
# DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════════
NAVY = RGBColor(0x1E, 0x27, 0x61)
DARK_NAVY = RGBColor(0x14, 0x1B, 0x45)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
MUTED = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
BORDER_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Calibri"


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════
def _shape(slide, shape_type, x, y, w, h, fill=None, line_color=None, line_width=None):
    """Add a shape with optional fill and line."""
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width or 1)
    else:
        s.line.fill.background()
    return s


def _text(slide, text, x, y, w, h, size=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.LEFT, font=FONT):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = str(text) if text else ""
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def _header(slide, title: str):
    """Add standard header bar."""
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.9, fill=NAVY)
    _text(slide, title, 0.8, 0.15, 10, 0.6, size=24, color=WHITE, bold=True)


def _footer(slide, project_name: str, page_num: int | None = None):
    """Add standard footer bar.

    When ``page_num`` is None, the footer derives the page number from
    the slide's actual position in its presentation. This makes the
    numbering robust to conditionally-added slides — for example,
    skipping the balance-sheet slide when no entity has balance data,
    which would otherwise leave a gap in the numbering (slide titled
    "page 9" but living at index 9, with "page 10" still labeled 10).

    Callers can still pass an explicit ``page_num`` for tests or
    when the position is known up-front, but the default path is the
    automatic derivation.
    """
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - 0.35, SLIDE_W, 0.35, fill=NAVY)
    _text(slide, f"{project_name}  |  Confidencial", 0.8, SLIDE_H - 0.32, 8, 0.3,
          size=9, color=ICE_BLUE)

    if page_num is None:
        # Derive from the slide's position in its presentation. The slide
        # was just added to the presentation, so it's the last one.
        prs = slide.part.package.presentation_part.presentation
        try:
            page_num = list(prs.slides).index(slide) + 1
        except ValueError:
            page_num = 0

    _text(slide, str(page_num), SLIDE_W - 1.5, SLIDE_H - 0.32, 0.7, 0.3,
          size=9, color=ICE_BLUE, align=PP_ALIGN.RIGHT)


def _card(slide, val, label, x, y, w=2.7, h=1.8, accent=ACCENT):
    """Add a KPI card with accent top bar."""
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.05, fill=accent)
    _text(slide, str(val), x + 0.25, y + 0.3, w - 0.5, 0.7, size=30, color=DARK_TEXT, bold=True)
    _text(slide, label, x + 0.25, y + 1.05, w - 0.5, 0.5, size=11, color=MUTED)


def _side_card(slide, val, label, sub, x, y, w=3.7, accent=ACCENT):
    """Add a metric card with left accent bar."""
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 1.5, fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, 1.5, fill=accent)
    _text(slide, str(val), x + 0.3, y + 0.15, w - 0.5, 0.6, size=26, color=DARK_TEXT, bold=True)
    _text(slide, label, x + 0.3, y + 0.75, w - 0.5, 0.3, size=11, color=MUTED)
    _text(slide, sub, x + 0.3, y + 1.05, w - 0.5, 0.25, size=11, color=accent, bold=True)


def _add_table(slide, headers, rows, x, y, w, row_h=0.35, highlight_row=-1):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows))
    tbl = tbl_shape.table

    col_w = w / n_cols
    for j in range(n_cols):
        tbl.columns[j].width = Inches(col_w)

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = str(h)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val) if val is not None else "—"
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = FONT
            if i == highlight_row:
                p.font.color.rgb = ACCENT
                p.font.bold = True
            else:
                p.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG

    return tbl_shape


def _chart_to_image(fig) -> io.BytesIO:
    """Convert matplotlib figure to PNG BytesIO."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor="#F5F7FA")
    plt.close(fig)
    buf.seek(0)
    return buf


def _fmt_num(val, unit=""):
    """Format number for display."""
    if val is None:
        return "—"
    if isinstance(val, float):
        if val >= 1_000:
            return f"{val:,.0f}{unit}"
        return f"{val:,.1f}{unit}"
    return f"{val}{unit}"


def _safe(val, default="—"):
    """Safely get a value or default."""
    if val is None:
        return default
    return str(val)


def _norm_label(s: str) -> str:
    """Lowercase + strip accents for accent-insensitive matching of DRE
    line labels.

    Brazilian financial packs are inconsistent about accents: the same
    workbook may have "Receita Líquida" on one sheet and "Receita
    Liquida" on another, depending on whether the row was typed by a
    human or auto-generated by a model. The pptx exporter previously
    matched ``"receita líquida" in label.lower()``, which is
    accent-sensitive and silently failed on Regenera (whose label has
    no accent). Same root cause as the E3.3 B7 fix in
    ``_extract_dre_value``; this is the analogous fix in the exporter.
    """
    if s is None:
        return ""
    return "".join(
        c for c in unicodedata.normalize("NFKD", s.lower())
        if not unicodedata.combining(c)
    )


# Common DRE labels we look for, normalized once. Match is "any token
# matches as a substring of the normalized label" — so "(=) Receita
# Liquida (RL)" still matches "receita liquida".
_REVENUE_LABEL_TOKENS = (
    "receita liquida",
    "(=) receita",
    "net revenue",
    "receita operacional liquida",
)
_EBITDA_LABEL_TOKENS = (
    "ebitda",
)
_EBITDA_EXCLUDE_TOKENS = (
    # Avoid matching "Margem EBITDA" or "EBITDA % RL" rows
    "margem",
    "%",
)


def _find_dre_line(stmt, token_list, exclude_tokens=()):
    """Return the first line in ``stmt.lines`` whose normalized label
    contains any token in ``token_list`` and contains none of
    ``exclude_tokens``. Returns None if no match.
    """
    if not stmt or not stmt.lines:
        return None
    for line in stmt.lines:
        norm = _norm_label(line.label)
        if not any(tok in norm for tok in token_list):
            continue
        if any(ex in norm for ex in exclude_tokens):
            continue
        return line
    return None


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════
def _slide_cover(prs, dossier: Dossier):
    """Slide 1: Dark cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_NAVY

    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.06, fill=ACCENT)

    project = dossier.metadata.project_name or "Projeto"
    company = dossier.metadata.target_company or "Empresa"

    _text(slide, project.upper(), 0.8, 1.8, 8, 0.5, size=16, color=ICE_BLUE)
    _text(slide, company, 0.8, 2.4, 8, 2, size=44, color=WHITE, bold=True)
    _text(slide, "Dossiê de Investimento", 0.8, 4.5, 8, 0.5, size=20, color=ICE_BLUE)

    # Confidential badge
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 5.6, 2.0, 0.35,
           fill=RGBColor(0x2D, 0x3A, 0x6E), line_color=RGBColor(0x4A, 0x5A, 0x9A), line_width=0.5)
    _text(slide, "CONFIDENCIAL", 0.8, 5.62, 2.0, 0.35, size=10, color=ICE_BLUE, bold=True, align=PP_ALIGN.CENTER)

    _MONTHS_PT = {1:"Janeiro",2:"Fevereiro",3:"Março",4:"Abril",5:"Maio",6:"Junho",
                  7:"Julho",8:"Agosto",9:"Setembro",10:"Outubro",11:"Novembro",12:"Dezembro"}
    now = datetime.now()
    _text(slide, f"{_MONTHS_PT[now.month]} {now.year}", 0.8, 6.2, 3, 0.4, size=12, color=MUTED)

    # Right side KPIs
    _shape(slide, MSO_SHAPE.RECTANGLE, 11.5, 1.5, 0.04, 4.5, fill=ACCENT)

    p = dossier.company.profile
    kpis = []
    if p.number_of_stores.value:
        kpis.append((str(p.number_of_stores.value), "Lojas"))
    if p.founding_year.value:
        kpis.append((str(p.founding_year.value), "Fundação"))

    for i, (val, label) in enumerate(kpis[:3]):
        y_pos = 1.8 + i * 1.5
        _text(slide, val, 11.8, y_pos, 1.3, 0.5, size=22, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
        _text(slide, label, 11.8, y_pos + 0.45, 1.3, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _slide_summary(prs, dossier: Dossier):
    """Slide 2: Executive summary with KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Sumário executivo")

    p = dossier.company.profile
    cards = []
    if p.number_of_stores.value:
        cards.append((str(p.number_of_stores.value), "Lojas em operação", ACCENT))
    if p.sector.value:
        cards.append((_safe(p.sector.value), "Setor", NAVY))
    if p.founding_year.value:
        cards.append((str(p.founding_year.value), "Ano de fundação", ACCENT_GREEN))

    card_w = 3.5
    gap = 0.35
    start_x = 0.8
    for i, (val, label, accent) in enumerate(cards[:4]):
        x = start_x + i * (card_w + gap)
        _card(slide, val, label, x, 1.4, w=card_w, accent=accent)

    # Description
    if p.description.value:
        _text(slide, "Sobre a empresa", 0.8, 3.7, 4, 0.35, size=16, color=NAVY, bold=True)
        _text(slide, str(p.description.value), 0.8, 4.15, 5.8, 1.5, size=12, color=DARK_TEXT)

    # Right side info
    _shape(slide, MSO_SHAPE.RECTANGLE, 7.2, 3.7, 0.04, 2.5, fill=ACCENT)

    info = []
    if p.business_model.value:
        info.append(("Modelo", str(p.business_model.value)[:80]))
    if p.headquarters.value:
        info.append(("Sede", str(p.headquarters.value)))
    if p.target_audience.value:
        info.append(("Público-Alvo", str(p.target_audience.value)[:80]))
    t = dossier.transaction
    if t.advisor.value:
        info.append(("Advisor", str(t.advisor.value)))

    for i, (label, value) in enumerate(info[:5]):
        y = 3.8 + i * 0.6
        _text(slide, label, 7.5, y, 2, 0.22, size=9, color=MUTED, bold=True)
        _text(slide, value, 7.5, y + 0.2, 5.5, 0.4, size=11, color=DARK_TEXT)

    _footer(slide, dossier.metadata.project_name)


def _slide_company(prs, dossier: Dossier):
    """Slide 3: Company profile."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "A empresa")
    p = dossier.company.profile

    fields = [
        ("Razão Social", _safe(p.legal_name.value)),
        ("Nome Fantasia", _safe(p.trade_name.value)),
        ("Setor", _safe(p.sector.value)),
        ("Modelo de Negócio", _safe(p.business_model.value)[:70]),
        ("Público-Alvo", _safe(p.target_audience.value)[:70]),
        ("Sede", _safe(p.headquarters.value)),
        ("Fundação", _safe(p.founding_year.value)),
        ("Nº Lojas", _safe(p.number_of_stores.value)),
        ("Nº Funcionários", _safe(p.number_of_employees.value)),
    ]

    # Two columns: long fields left, short fields right
    col1 = fields[:5]
    col2 = fields[5:]

    for i, (label, val) in enumerate(col1):
        y = 1.3 + i * 0.75
        _text(slide, label, 0.8, y, 2.5, 0.22, size=10, color=MUTED, bold=True)
        _text(slide, val, 0.8, y + 0.2, 5.5, 0.45, size=12, color=DARK_TEXT)
        if i < len(col1) - 1:
            _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, y + 0.65, 5.5, 0.007, fill=BORDER_LIGHT)

    for i, (label, val) in enumerate(col2):
        y = 1.3 + i * 0.75
        _text(slide, label, 7.3, y, 2.5, 0.22, size=10, color=MUTED, bold=True)
        _text(slide, val, 7.3, y + 0.2, 5.0, 0.45, size=12, color=DARK_TEXT)
        if i < len(col2) - 1:
            _shape(slide, MSO_SHAPE.RECTANGLE, 7.3, y + 0.65, 5.0, 0.007, fill=BORDER_LIGHT)

    # Description box
    if p.description.value:
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 5.2, 11.7, 1.3,
               fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 5.2, 0.06, 1.3, fill=ACCENT)
        _text(slide, "Descrição", 1.1, 5.3, 3, 0.3, size=12, color=NAVY, bold=True)
        _text(slide, str(p.description.value), 1.1, 5.6, 11.0, 0.8, size=11, color=DARK_TEXT)

    _footer(slide, dossier.metadata.project_name)


def _slide_executives(prs, dossier: Dossier):
    """Slide 4: Directors and shareholders table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Diretoria e acionistas")

    execs = dossier.company.executives
    if execs:
        headers = ["Nome", "Cargo", "Participação", "Background"]
        rows = []
        for ex in execs:
            pct = f"{ex.ownership_pct}%" if ex.ownership_pct else "—"
            bg = (ex.background or "—")[:80]
            rows.append([ex.name, ex.role or "—", pct, bg])

        _add_table(slide, headers, rows, 0.8, 1.3, 11.7)
    else:
        _text(slide, "Dados de diretoria não disponíveis", 0.8, 2, 8, 0.5, size=14, color=MUTED)

    _footer(slide, dossier.metadata.project_name)


def _slide_timeline(prs, dossier: Dossier):
    """Slide 5: Company timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Histórico")

    events = sorted(dossier.company.timeline, key=lambda e: e.year)

    if not events:
        _text(slide, "Timeline não disponível", 0.8, 2, 8, 0.5, size=14, color=MUTED)
        _footer(slide, dossier.metadata.project_name)
        return

    # Timeline as horizontal steps
    n = len(events)
    usable_w = 11.0
    step_w = min(usable_w / n, 2.5)
    start_x = 0.8 + (usable_w - step_w * n) / 2

    # Horizontal line
    line_y = 2.8
    _shape(slide, MSO_SHAPE.RECTANGLE, start_x, line_y, step_w * n, 0.03, fill=NAVY)

    for i, ev in enumerate(events):
        cx = start_x + i * step_w + step_w / 2

        # Dot on the line
        dot_size = 0.18
        _shape(slide, MSO_SHAPE.OVAL, cx - dot_size / 2, line_y - dot_size / 2 + 0.015,
               dot_size, dot_size, fill=ACCENT)

        # Year above
        _text(slide, str(ev.year), cx - 0.6, line_y - 0.55, 1.2, 0.4,
              size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

        # Description below
        desc = ev.description[:60] if len(ev.description) > 60 else ev.description
        _text(slide, desc, cx - step_w / 2, line_y + 0.3, step_w - 0.1, 1.2,
              size=9, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    _footer(slide, dossier.metadata.project_name)


def _slide_products(prs, dossier: Dossier):
    """Slide 6: Products and brands."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Produtos e marcas")

    products = dossier.company.products
    if not products:
        _text(slide, "Dados de produtos não disponíveis", 0.8, 2, 8, 0.5, size=14, color=MUTED)
        _footer(slide, dossier.metadata.project_name)
        return

    # Split into categories and proprietary brands
    categories = [p for p in products if not p.is_proprietary]
    brands = [p for p in products if p.is_proprietary]

    # Categories table
    if categories:
        _text(slide, "Categorias de produto", 0.8, 1.2, 5, 0.35, size=14, color=NAVY, bold=True)
        headers = ["Produto", "Categoria", "% Receita"]
        rows = [[p.name, p.category or "—", f"{p.revenue_share_pct}%" if p.revenue_share_pct else "—"]
                for p in categories]
        _add_table(slide, headers, rows, 0.8, 1.6, 5.5)

    # Brands
    if brands:
        _text(slide, "Marcas próprias", 7.0, 1.2, 5, 0.35, size=14, color=NAVY, bold=True)
        headers = ["Marca", "Categoria"]
        rows = [[p.name, p.category or "—"] for p in brands]
        _add_table(slide, headers, rows, 7.0, 1.6, 5.5)

    _footer(slide, dossier.metadata.project_name)


def _slide_dre(prs, dossier: Dossier):
    """Slide 7: Income statement (DRE) — chart of largest entity + a
    summary table showing Receita Líquida and EBITDA per entity, last 3
    years, in a compact grid.

    Previously this slide rendered a chart for one entity plus a list
    of entity names with "DRE: N anos" labels — visually empty for
    multi-entity CIMs like Regenera (6 entities). The chart silently
    disappeared whenever the receita-líquida label had no accent,
    leaving only the labels behind.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Demonstração de resultados (DRE)")

    fin = dossier.financials
    available = [(e.name, e.dre) for e in fin.entities if e.dre and e.dre.lines]

    if not available:
        _text(slide, "Dados financeiros não disponíveis", 0.8, 2, 8, 0.5,
              size=14, color=MUTED)
        _footer(slide, dossier.metadata.project_name)
        return

    # ── Pick the entity for the chart: largest revenue in the most
    #    recent year wins. Falls back to the first entity if matching
    #    fails on every entity.
    chart_pick = None
    chart_pick_revenue = -1.0
    for name, stmt in available:
        rev_line = _find_dre_line(stmt, _REVENUE_LABEL_TOKENS)
        if not rev_line or not stmt.years:
            continue
        last_y = stmt.years[-1]
        v = rev_line.values.get(last_y) or 0
        if v > chart_pick_revenue:
            chart_pick_revenue = v
            chart_pick = (name, stmt, rev_line)

    # ── Chart panel (left) ─────────────────────────────────────────────
    if chart_pick:
        name, stmt, receita_line = chart_pick
        years = stmt.years or []
        ebitda_line = _find_dre_line(stmt, _EBITDA_LABEL_TOKENS,
                                     exclude_tokens=_EBITDA_EXCLUDE_TOKENS)

        fig, ax = plt.subplots(figsize=(7, 3.5))
        fig.patch.set_facecolor("#F5F7FA")
        ax.set_facecolor("#F5F7FA")

        vals_receita = [receita_line.values.get(y, 0) or 0 for y in years]
        x = range(len(years))
        bars = ax.bar(x, vals_receita, width=0.6, color="#1E2761",
                      label="Receita Líquida", zorder=3)

        if vals_receita and max(vals_receita) > 0:
            for bar, val in zip(bars, vals_receita):
                if val:
                    ax.text(bar.get_x() + bar.get_width() / 2,
                            bar.get_height() + max(vals_receita) * 0.02,
                            f"{val:,.0f}", ha="center", va="bottom",
                            fontsize=7, color="#1E2761", fontweight="bold")

        if ebitda_line:
            vals_ebitda = [ebitda_line.values.get(y, 0) or 0 for y in years]
            ax.plot(x, vals_ebitda, color="#3B82F6", linewidth=2.5,
                    marker="o", markersize=5, label="EBITDA", zorder=4)
            if vals_receita and max(vals_receita) > 0:
                for i, val in enumerate(vals_ebitda):
                    if val:
                        ax.text(i, val + max(vals_receita) * 0.03,
                                f"{val:,.0f}", ha="center", fontsize=7,
                                color="#3B82F6")

        ax.set_xticks(x)
        ax.set_xticklabels(years, fontsize=8, color="#64748B")
        ax.set_ylabel("BRL k", fontsize=9, color="#64748B")
        ax.legend(fontsize=8, frameon=False)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#E2E8F0")
        ax.spines["bottom"].set_color("#E2E8F0")
        ax.tick_params(colors="#64748B", labelsize=7)
        ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)
        plt.tight_layout()

        img = _chart_to_image(fig)
        slide.shapes.add_picture(img, Inches(0.5), Inches(1.1),
                                 Inches(7.5), Inches(3.8))
        _text(slide, f"Gráfico: {name}", 0.8, 5.0, 4, 0.3, size=9, color=MUTED)
    else:
        # No matchable revenue line on any entity. Don't silently render
        # a "Gráfico:" label with nothing behind it.
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.5, 1.1, 7.5, 3.8,
               fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
        _text(slide, "Linha de Receita Líquida não identificada — verifique os dados",
              0.8, 2.5, 7, 0.5, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Summary table (right): Receita Líquida + EBITDA per entity,
    #    last 3 years. With 6 entities × 2 metrics × 3 years it fits
    #    comfortably. ───────────────────────────────────────────────────
    if chart_pick:
        years_for_table = (chart_pick[1].years or [])[-3:]   # last 3
    else:
        # Fallback: use the first entity's years
        years_for_table = (available[0][1].years or [])[-3:]

    headers = ["Entidade", "Métrica"] + list(years_for_table)
    rows = []
    for name, stmt in available:
        rev_line = _find_dre_line(stmt, _REVENUE_LABEL_TOKENS)
        ebt_line = _find_dre_line(stmt, _EBITDA_LABEL_TOKENS,
                                  exclude_tokens=_EBITDA_EXCLUDE_TOKENS)
        # Two rows per entity: revenue and EBITDA
        if rev_line:
            row = [name, "Rec. Líq."] + [
                _fmt_num(rev_line.values.get(y)) for y in years_for_table
            ]
        else:
            row = [name, "Rec. Líq."] + ["—" for _ in years_for_table]
        rows.append(row)
        if ebt_line:
            row = ["", "EBITDA"] + [
                _fmt_num(ebt_line.values.get(y)) for y in years_for_table
            ]
        else:
            row = ["", "EBITDA"] + ["—" for _ in years_for_table]
        rows.append(row)

    # Table on the right side; height auto-grows with row count
    n_rows = len(rows) + 1   # +1 for header
    row_h = 0.22 if n_rows > 8 else 0.28
    _add_table(slide, headers, rows,
               x=8.4, y=1.1, w=4.7, row_h=row_h)

    _footer(slide, dossier.metadata.project_name)


def _slide_cash_flow(prs, dossier: Dossier):
    """Slide 8: Cash flow placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Fluxo de caixa")

    _shape(slide, MSO_SHAPE.RECTANGLE, 2.5, 2.5, 8.3, 2.5, fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
    _shape(slide, MSO_SHAPE.RECTANGLE, 2.5, 2.5, 8.3, 0.05, fill=ACCENT_AMBER)
    _text(slide, "Fluxo de caixa ainda não disponível", 2.5, 3.0, 8.3, 0.5,
          size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, "Os dados de fluxo de caixa (operacional, investimento, financiamento) "
                 "serão incluídos quando disponíveis no documento fonte ou inputados manualmente.",
          3.0, 3.6, 7.3, 1.0, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    _footer(slide, dossier.metadata.project_name)


def _slide_balance(prs, dossier: Dossier):
    """Slide 9: Balance sheet summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Balanço patrimonial")

    fin = dossier.financials
    stmts = [(e.name, e.balance_sheet) for e in fin.entities]

    available = [(name, stmt) for name, stmt in stmts if stmt and stmt.lines]

    if not available:
        _text(slide, "Dados de balanço não disponíveis", 0.8, 2, 8, 0.5, size=14, color=MUTED)
        _footer(slide, dossier.metadata.project_name)
        return

    # Show summary for each entity
    for idx, (ent_name, stmt) in enumerate(available[:3]):
        x_start = 0.8 + idx * 4.1
        _text(slide, ent_name, x_start, 1.2, 3.7, 0.35, size=14, color=NAVY, bold=True)

        # Show key lines
        key_labels = ["ativo total", "passivo total", "patrimônio líquido"]
        y_pos = 1.6
        for line in stmt.lines:
            label_lower = line.label.lower()
            if any(kl in label_lower for kl in key_labels):
                last_year = stmt.years[-1] if stmt.years else None
                val = line.values.get(last_year) if last_year else None
                _text(slide, line.label[:30], x_start, y_pos, 2.5, 0.25, size=10, color=MUTED)
                _text(slide, _fmt_num(val), x_start + 2.5, y_pos, 1.2, 0.25,
                      size=11, color=DARK_TEXT, bold=True, align=PP_ALIGN.RIGHT)
                y_pos += 0.35

    _footer(slide, dossier.metadata.project_name)


def _slide_market(prs, dossier: Dossier):
    """Slide 10: Market size and growth.

    The slide always shows three rows: market-size cards, fragmentation
    note, and growth drivers. When any of those isn't populated in the
    dossier, we still render the section header with a "—" placeholder
    so the reader sees what's intentionally missing rather than a blank
    region. (Previously, sections silently disappeared, leaving the
    slide visually sparse on CIMs that lacked market detail.)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Mercado")

    # Market size cards
    ms_list = dossier.market.market_sizes
    if ms_list:
        for i, ms in enumerate(ms_list[:3]):
            x = 0.8 + i * 4.1
            val = f"{ms.unit} {ms.value:,.1f}" if ms.value else f"{ms.unit} —"
            label = f"Mercado {ms.geography} {ms.year}"
            # Auto-detect CAGR format: <1 = decimal (0.033), >=1 = already percentage (3.3)
            if ms.cagr:
                cagr_pct = ms.cagr * 100 if ms.cagr < 1 else ms.cagr
                cagr = f"CAGR {cagr_pct:.1f}%"
            else:
                cagr = "—"
            accent = [ACCENT, ACCENT_GREEN, NAVY][i % 3]
            _side_card(slide, val, label, cagr, x, 1.3, accent=accent)
    else:
        # Empty placeholder card so the layout doesn't look broken
        _side_card(slide, "—", "Tamanho de mercado", "Não extraído",
                   0.8, 1.3, accent=MUTED)
        _text(slide, "Sugestão: Sebrae, IBGE, relatórios setoriais",
              4.7, 1.5, 7, 0.3, size=10, color=MUTED)

    # Fragmentation
    _text(slide, "Fragmentação", 0.8, 3.3, 4, 0.3, size=14, color=NAVY, bold=True)
    if dossier.market.market_fragmentation.is_filled:
        _text(slide, str(dossier.market.market_fragmentation.value), 0.8, 3.65, 11.5, 0.6,
              size=11, color=DARK_TEXT)
    else:
        _text(slide, "—  (estrutura de mercado não identificada)",
              0.8, 3.65, 11.5, 0.6, size=11, color=MUTED)

    # Growth drivers
    _text(slide, "Drivers de crescimento", 0.8, 4.4, 4, 0.3, size=14, color=NAVY, bold=True)
    drivers = dossier.market.growth_drivers
    if drivers:
        driver_text = "; ".join(str(d.value) for d in drivers[:5])
        _text(slide, driver_text, 0.8, 4.75, 11.5, 0.6, size=11, color=DARK_TEXT)
    else:
        _text(slide, "—  (drivers de crescimento não identificados)",
              0.8, 4.75, 11.5, 0.6, size=11, color=MUTED)

    _footer(slide, dossier.metadata.project_name)


def _slide_competitors(prs, dossier: Dossier):
    """Slide 11: Competitive landscape.

    Renders three sections — competitor table, sector multiples,
    precedent transactions — each with an explicit placeholder when
    the dossier has nothing to show. Previously a Regenera-style
    dossier with 0 competitors and 0 precedent transactions left
    the entire left-hand pane empty; the multiples cards on the
    right alone made the slide look broken.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Landscape competitivo")

    # Left side: competitors
    comps = dossier.market.competitors
    _text(slide, "Concorrentes", 0.8, 1.0, 4, 0.3, size=14, color=NAVY, bold=True)
    if comps:
        # Find our company's position for highlighting
        company_name = (dossier.metadata.target_company or "").lower()
        highlight = -1
        for i, c in enumerate(comps):
            if company_name and company_name in c.name.lower():
                highlight = i
                break

        headers = ["Empresa", "Lojas", "Receita", "Investidor"]
        rows = [[c.name, _safe(c.stores), _safe(c.revenue),
                 _safe(c.investor)] for c in comps]
        _add_table(slide, headers, rows, 0.8, 1.35, 7.5, highlight_row=highlight)
    else:
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 1.35, 7.5, 1.5,
               fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
        _text(slide, "Concorrentes não identificados", 0.8, 1.7, 7.5, 0.4,
              size=12, color=MUTED, align=PP_ALIGN.CENTER)
        _text(slide, "Sugestão: pesquisa setorial / sites das principais empresas",
              0.8, 2.1, 7.5, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Multiples (right side)
    if dossier.market.global_multiples_median.is_filled:
        multiples = dossier.market.global_multiples_median.value
        if isinstance(multiples, dict):
            _text(slide, "Múltiplos (mediana setorial)", 9.0, 1.0, 3.8, 0.3,
                  size=14, color=NAVY, bold=True)

            ev_rev = multiples.get("ev_revenue_median")
            ev_ebitda = multiples.get("ev_ebitda_median")

            if ev_rev:
                _shape(slide, MSO_SHAPE.RECTANGLE, 9.0, 1.35, 3.8, 0.8,
                       fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
                _text(slide, f"{ev_rev}x", 9.2, 1.4, 2, 0.4, size=24,
                      color=DARK_TEXT, bold=True)
                _text(slide, "EV/Receita", 9.2, 1.8, 2, 0.25, size=10, color=MUTED)

            if ev_ebitda:
                _shape(slide, MSO_SHAPE.RECTANGLE, 9.0, 2.3, 3.8, 0.8,
                       fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
                _text(slide, f"{ev_ebitda}x", 9.2, 2.35, 2, 0.4, size=24,
                      color=DARK_TEXT, bold=True)
                _text(slide, "EV/EBITDA", 9.2, 2.75, 2, 0.25, size=10, color=MUTED)

            # Provenance line — analysts will care whether multiples
            # came from CIM, web search, or manual override
            note = multiples.get("source_note") or ""
            if note:
                _text(slide, note[:60], 9.0, 3.2, 3.8, 0.3,
                      size=8, color=MUTED)

    # Precedent transactions
    txns = dossier.market.precedent_transactions
    _text(slide, "Transações precedentes", 0.8, 4.0, 6, 0.3,
          size=14, color=NAVY, bold=True)
    if txns:
        headers = ["Comprador", "Alvo", "Valor", "EV/EBITDA"]
        rows = [[t.buyer, t.target, _safe(t.value),
                 f"{t.ev_ebitda}x" if t.ev_ebitda else "—"]
                for t in txns[:6]]
        _add_table(slide, headers, rows, 0.8, 4.35, 11.7)
    else:
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 4.35, 11.7, 1.0,
               fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
        _text(slide, "Transações precedentes não identificadas",
              0.8, 4.55, 11.7, 0.3, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        _text(slide, "Sugestão: PitchBook, Mergermarket, releases dos compradores",
              0.8, 4.85, 11.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.CENTER)

    _footer(slide, dossier.metadata.project_name)


def _slide_transaction(prs, dossier: Dossier):
    """Slide 12: Transaction / deal structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Oportunidade de investimento")

    t = dossier.transaction
    fields = [
        ("Tipo", _safe(t.transaction_type.value)),
        ("Contexto", _safe(t.context.value)),
        ("Stake alvo", _safe(t.target_stake_range.value)),
        ("Capital necessário", _safe(t.capital_needed.value)),
        ("Uso dos recursos", _safe(t.use_of_proceeds.value)),
        ("Perímetro", _safe(t.perimeter.value)),
        ("Advisor", _safe(t.advisor.value)),
    ]

    for i, (label, val) in enumerate(fields):
        y = 1.3 + i * 0.7
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, y, 11.7, 0.6, fill=WHITE if i % 2 == 0 else LIGHT_BG,
               line_color=BORDER_LIGHT, line_width=0.3)
        _text(slide, label, 1.0, y + 0.1, 2.5, 0.35, size=11, color=NAVY, bold=True)
        _text(slide, val, 3.5, y + 0.1, 8.5, 0.35, size=11, color=DARK_TEXT)

    _footer(slide, dossier.metadata.project_name)


def _slide_gaps(prs, dossier: Dossier):
    """Slide 13: Gap analysis and next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Gaps e próximos passos")

    gaps = dossier.gaps
    critical = [g for g in gaps if g.severity == "critical"]
    important = [g for g in gaps if g.severity == "important"]

    # Critical gaps
    if critical:
        _text(slide, f"Gaps críticos ({len(critical)})", 0.8, 1.2, 5, 0.35, size=14, color=ACCENT_RED, bold=True)
        for i, g in enumerate(critical[:5]):
            y = 1.65 + i * 0.4
            _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, y, 0.06, 0.3, fill=ACCENT_RED)
            _text(slide, g.description, 1.1, y, 5.5, 0.3, size=10, color=DARK_TEXT)
            if g.suggested_source:
                _text(slide, f"→ {g.suggested_source}", 4.5, y, 3, 0.3, size=9, color=MUTED)

    # Important gaps
    if important:
        y_start = 1.65 + min(len(critical), 5) * 0.4 + 0.3
        _text(slide, f"Gaps importantes ({len(important)})", 0.8, y_start, 5, 0.35, size=14, color=ACCENT_AMBER, bold=True)
        for i, g in enumerate(important[:6]):
            y = y_start + 0.4 + i * 0.35
            _shape(slide, MSO_SHAPE.RECTANGLE, 0.8, y, 0.06, 0.25, fill=ACCENT_AMBER)
            _text(slide, g.description, 1.1, y, 5.5, 0.25, size=10, color=DARK_TEXT)

    # Summary card
    _shape(slide, MSO_SHAPE.RECTANGLE, 8.5, 1.2, 4.0, 2.0, fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
    _text(slide, "Resumo", 8.8, 1.35, 3, 0.3, size=14, color=NAVY, bold=True)
    _text(slide, f"Total de gaps: {len(gaps)}", 8.8, 1.8, 3, 0.25, size=12, color=DARK_TEXT)
    _text(slide, f"Críticos: {len(critical)}", 8.8, 2.1, 3, 0.25, size=12, color=ACCENT_RED, bold=True)
    _text(slide, f"Importantes: {len(important)}", 8.8, 2.4, 3, 0.25, size=12, color=ACCENT_AMBER)

    # Web gaps
    web_gaps = [g for g in gaps if g.requires_internet]
    if web_gaps:
        _text(slide, f"Requerem pesquisa web: {len(web_gaps)}", 8.8, 2.7, 3, 0.25, size=10, color=MUTED)

    _footer(slide, dossier.metadata.project_name)


# ═══════════════════════════════════════════════════════════════
# VALUATION SLIDES
# ═══════════════════════════════════════════════════════════════
def _slide_valuation_table(prs, dossier: Dossier, valuation_data: dict):
    """Slide 14: Valuation comparison table with KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Valuation — Cenários")

    summaries = valuation_data.get("summaries", [])
    inputs = valuation_data.get("inputs", {})

    if not summaries:
        _text(slide, "Dados de valuation não disponíveis", 0.8, 2, 8, 0.5, size=14, color=MUTED)
        _footer(slide, dossier.metadata.project_name)
        return

    # KPI cards: WACC, Stake, IRR (base), MOIC (base)
    wacc = inputs.get("wacc", {})
    base = next((s for s in summaries if s.get("scenario_name", "").lower() == "base"), summaries[0])

    cards = [
        (f"{wacc.get('wacc', 0)*100:.1f}%", "WACC", NAVY),
        (f"{inputs.get('stake_pct', 0)*100:.0f}%", "Stake investidor", ACCENT),
        (f"{base.get('irr', 0)*100:.1f}%", "IRR (caso base)", ACCENT_GREEN),
        (f"{base.get('moic', 0):.2f}x", "MOIC (caso base)", ACCENT_AMBER),
    ]

    card_w = 2.8
    gap = 0.3
    start_x = 0.8
    for i, (val, label, accent) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        _card(slide, val, label, x, 1.2, w=card_w, h=1.5, accent=accent)

    # Comparison table
    headers = ["Cenário", "DCF Perp", "DCF Exit", "EV/EBITDA", "EV/Rev", "IRR", "MOIC"]
    rows = []
    base_row = -1
    for i, s in enumerate(summaries):
        name = s.get("scenario_name", f"Cenário {i+1}")
        if name.lower() == "base":
            base_row = i
        rows.append([
            name,
            f"{s.get('dcf_perpetuity', 0):,.0f}",
            f"{s.get('dcf_exit_multiple', 0):,.0f}",
            f"{s.get('multiples_ev_ebitda', 0):,.0f}",
            f"{s.get('multiples_ev_revenue', 0):,.0f}",
            f"{s.get('irr', 0)*100:.1f}%",
            f"{s.get('moic', 0):.2f}x",
        ])

    _add_table(slide, headers, rows, 0.8, 3.1, 11.7, row_h=0.4, highlight_row=base_row)

    # Unit note
    _text(slide, "EV em BRL k  |  Entry: DCF equity  |  Exit: múltiplo × EBITDA terminal",
          0.8, 5.0, 8, 0.25, size=9, color=MUTED)

    # WACC breakdown on right
    _shape(slide, MSO_SHAPE.RECTANGLE, 9.5, 5.3, 3.0, 1.5,
           fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
    _text(slide, "WACC", 9.7, 5.35, 2, 0.25, size=11, color=NAVY, bold=True)
    wacc_lines = [
        f"Ke: {wacc.get('cost_of_equity', 0)*100:.1f}%",
        f"Kd pós-tax: {wacc.get('cost_of_debt_aftertax', 0)*100:.1f}%",
        f"E/V: {wacc.get('equity_to_total', 0)*100:.0f}% | D/V: {wacc.get('debt_to_total', 0)*100:.0f}%",
    ]
    for j, line in enumerate(wacc_lines):
        _text(slide, line, 9.7, 5.65 + j * 0.3, 2.6, 0.25, size=9, color=DARK_TEXT)

    _footer(slide, dossier.metadata.project_name)


def _slide_valuation_chart(prs, dossier: Dossier, valuation_data: dict):
    """Slide 15: EV range chart (horizontal bars per scenario)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    _header(slide, "Valuation — Faixa de EV")

    summaries = valuation_data.get("summaries", [])
    if not summaries:
        _text(slide, "Dados de valuation não disponíveis", 0.8, 2, 8, 0.5, size=14, color=MUTED)
        _footer(slide, dossier.metadata.project_name)
        return

    # Build horizontal bar chart showing EV range per scenario
    fig, ax = plt.subplots(figsize=(8, 4))
    fig.patch.set_facecolor("#F5F7FA")
    ax.set_facecolor("#F5F7FA")

    names = [s.get("scenario_name", "?") for s in summaries]
    lows = [s.get("equity_range_low", 0) / 1000 for s in summaries]  # Convert to MM
    highs = [s.get("equity_range_high", 0) / 1000 for s in summaries]
    dcf_perps = [s.get("dcf_perpetuity", 0) / 1000 for s in summaries]
    dcf_exits = [s.get("dcf_exit_multiple", 0) / 1000 for s in summaries]

    y_pos = range(len(names))
    colors_bar = ["#EF4444", "#1E2761", "#10B981"]  # red, navy, green
    colors_bar = colors_bar[:len(names)]

    # Draw range bars
    for i, (name, lo, hi) in enumerate(zip(names, lows, highs)):
        # Range bar (low to high)
        ax.barh(i, hi - lo, left=lo, height=0.5, color=colors_bar[i], alpha=0.3, zorder=2)
        # DCF perpetuity marker
        ax.plot(dcf_perps[i], i, 'D', color=colors_bar[i], markersize=10, zorder=4)
        # DCF exit marker
        ax.plot(dcf_exits[i], i, 's', color=colors_bar[i], markersize=9, zorder=4, alpha=0.7)
        # Labels
        ax.text(lo - max(highs) * 0.02, i, f"{lo:.0f}",
                ha="right", va="center", fontsize=8, color=colors_bar[i], fontweight="bold")
        ax.text(hi + max(highs) * 0.02, i, f"{hi:.0f}",
                ha="left", va="center", fontsize=8, color=colors_bar[i], fontweight="bold")

    ax.set_yticks(y_pos)
    ax.set_yticklabels(names, fontsize=11, color="#1E2761", fontweight="bold")
    ax.set_xlabel("Enterprise Value (BRL MM)", fontsize=10, color="#64748B")
    ax.invert_yaxis()

    # Legend
    from matplotlib.lines import Line2D
    legend_elements = [
        Line2D([0], [0], marker='D', color='w', markerfacecolor='#64748B', markersize=8, label='DCF Perpetuidade'),
        Line2D([0], [0], marker='s', color='w', markerfacecolor='#64748B', markersize=8, label='DCF Exit Multiple'),
    ]
    ax.legend(handles=legend_elements, loc='lower right', fontsize=8, frameon=False)

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#E2E8F0")
    ax.spines["bottom"].set_color("#E2E8F0")
    ax.tick_params(colors="#64748B", labelsize=8)
    ax.xaxis.grid(True, color="#E2E8F0", linewidth=0.5, zorder=0)
    ax.set_axisbelow(True)
    plt.tight_layout()

    img = _chart_to_image(fig)
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.1), Inches(8.5), Inches(4.5))

    # Right side: "What needs to be true" for base case
    scenarios_data = valuation_data.get("scenarios", {})
    base_sc = scenarios_data.get("base", {})
    wnbt = base_sc.get("what_needs_to_be_true", [])

    if wnbt:
        _shape(slide, MSO_SHAPE.RECTANGLE, 9.3, 1.1, 3.5, 4.5,
               fill=WHITE, line_color=BORDER_LIGHT, line_width=0.5)
        _shape(slide, MSO_SHAPE.RECTANGLE, 9.3, 1.1, 0.06, 4.5, fill=ACCENT)
        _text(slide, "What needs to be true", 9.6, 1.2, 3.0, 0.3, size=12, color=NAVY, bold=True)
        _text(slide, "(caso base)", 9.6, 1.45, 3.0, 0.2, size=9, color=MUTED)

        for j, item in enumerate(wnbt[:6]):
            y = 1.8 + j * 0.6
            cat = item.get("category", "")
            cond = item.get("condition", "")
            _text(slide, cat, 9.6, y, 3.0, 0.2, size=9, color=ACCENT, bold=True)
            _text(slide, cond[:55], 9.6, y + 0.2, 3.0, 0.3, size=9, color=DARK_TEXT)

    # Terminal metrics
    _text(slide, "Faixa: DCF Perpetuidade (◆) a Múltiplos (barra) | BRL MM",
          0.8, 5.8, 8, 0.25, size=9, color=MUTED)

    _footer(slide, dossier.metadata.project_name)


# ═══════════════════════════════════════════════════════════════
# MAIN EXPORT FUNCTION
# ═══════════════════════════════════════════════════════════════
def export_pptx(dossier: Dossier, output_path: str, valuation_data: dict | None = None,
                verbose: bool = False) -> int:
    """Export dossier to a formatted PowerPoint presentation.

    Args:
        dossier: The dossier to export
        output_path: Path for the output .pptx file
        valuation_data: Optional valuation results from run_full_valuation
        verbose: Print progress

    Returns:
        Number of slides created
    """
    if verbose:
        print("  [PPT] Gerando apresentação...")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Build all slides
    _slide_cover(prs, dossier)
    _slide_summary(prs, dossier)
    _slide_company(prs, dossier)
    _slide_executives(prs, dossier)
    _slide_timeline(prs, dossier)
    _slide_products(prs, dossier)
    _slide_dre(prs, dossier)
    _slide_cash_flow(prs, dossier)

    # Balance sheet slide is added only when at least one entity has
    # extracted balance lines. Most pitch-deck CIMs (and financial
    # packs that ship just DREs, like Regenera) don't include balance
    # data, and an empty "Dados não disponíveis" slide adds noise to
    # the deck. The omission is also reflected in the gaps section,
    # so the absence is still visible to the reader.
    has_balance = any(
        e.balance_sheet and e.balance_sheet.lines
        for e in dossier.financials.entities
    )
    if has_balance:
        _slide_balance(prs, dossier)
    _slide_market(prs, dossier)
    _slide_competitors(prs, dossier)
    _slide_transaction(prs, dossier)

    # Valuation slides (14-15) — only if data available
    if valuation_data and valuation_data.get("summaries"):
        _slide_valuation_table(prs, dossier, valuation_data)
        _slide_valuation_chart(prs, dossier, valuation_data)

    _slide_gaps(prs, dossier)  # Always last

    # Save
    n_slides = len(prs.slides)
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    if verbose:
        print(f"  [PPT] ✅ Salvo em: {output_path} ({n_slides} slides)")

    return n_slides