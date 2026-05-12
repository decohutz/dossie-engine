"""
Tests for E5 — visual quality of XLSX and PPTX outputs.

The four problems addressed:

* **P1** — slide 7 ("Demonstração de resultados") rendered for Regenera
  v008 with no chart and no table, just labels listing "DRE: N anos"
  for each entity. Two root causes: (a) the receita-líquida line
  matcher was accent-sensitive ("receita líquida" in label.lower())
  and missed Regenera's "Receita Liquida" without accent; (b) when
  the matcher returned None, the chart was silently skipped without
  fallback. Now: accent-insensitive matcher (same fix as B7 from
  E3.3), fallback message when no revenue line is found, and a
  multi-entity summary table replacing the per-entity label list.

* **P6** — the balance-sheet slide rendered "Dados de balanço não
  disponíveis" as a full slide on Regenera (which has no balance
  data). Slide is now skipped entirely when no entity has balance
  lines, and footer page numbers auto-derive from the slide's
  position so renumbering stays consistent.

* **P7** — the "Visão Geral" sheet left cells visually blank when
  source values were None (founding year, employee count, ownership
  pct). Now: missing values render as "—" in all cases, with the
  proper number format applied only when the value is non-null.

* **P5** — the Mercado and Landscape slides looked sparse on CIMs
  with thin market data. Now: empty sections render an explicit
  placeholder ("—" or "X não identificados") so the layout stays
  consistent across CIMs with very different market-data depth.
"""
from __future__ import annotations

import pytest

openpyxl = pytest.importorskip("openpyxl")
pptx_module = pytest.importorskip("pptx")

from tests.test_xlsx_parser import _new_workbook, _add_cover, _add_dre_sheet


def _make_dossier_with_2_entities(tmp_path):
    """Build a 2-entity Regenera-shaped dossier from a small XLSX.
    Years intentionally have NO 'E' suffix so the parser strips them
    and the consolidated stays projected, exercising the P1 path."""
    wb = _new_workbook()
    _add_cover(wb, last_update_year=2025)
    rows_a = [
        ("Receita Bruta", 100_000, 110_000, 130_000, 150_000, 175_000, 200_000),
        ("Impostos",       -15_000, -16_500, -19_500, -22_500, -26_250, -30_000),
        # NOTE: deliberately no accent on "Liquida" — exercises the P1.a fix
        ("Receita Liquida", 85_000,  93_500, 110_500, 127_500, 148_750, 170_000),
        ("CMV",            -30_000, -32_000, -38_000, -44_000, -51_000, -58_000),
        ("EBITDA",          30_000,  34_500,  40_500,  46_500,  54_750,  64_000),
    ]
    rows_b = [
        ("Receita Bruta", 50_000, 60_000, 75_000, 90_000, 105_000, 125_000),
        ("Receita Liquida", 42_500, 51_000, 63_750, 76_500, 89_250, 106_250),
        ("EBITDA",         15_000, 18_000, 23_000, 28_000, 33_500, 40_000),
    ]
    _add_dre_sheet(wb, "DRE Alpha", rows=rows_a)
    _add_dre_sheet(wb, "DRE Beta",  rows=rows_b)
    path = tmp_path / "p.xlsx"
    wb.save(path)
    return path


# ── P1 — Slide DRE ──────────────────────────────────────────────────
def test_dre_slide_renders_chart_for_accent_free_label(tmp_path):
    """Regression: 'Receita Liquida' (no accent) must produce a chart.
    Before P1.a, the accent-sensitive matcher silently dropped the chart.
    """
    from src.pipeline.orchestrator import run_pipeline
    from src.exporters.pptx_exporter import export_pptx
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    path = _make_dossier_with_2_entities(tmp_path)
    dossier = run_pipeline(inputs=[str(path)], use_llm=False, project_name="t")

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    prs = Presentation(str(out))
    dre_slide = prs.slides[6]   # 0-indexed slide 7

    # Must have at least one PICTURE shape (the chart)
    has_picture = any(
        s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in dre_slide.shapes
    )
    assert has_picture, "DRE slide should contain a chart picture"


def test_dre_slide_has_summary_table_with_all_entities(tmp_path):
    """The new DRE slide includes a summary table covering every
    operating entity, not just the first one. That replaces the old
    'list of entity labels' panel that made multi-entity decks look
    empty."""
    from src.pipeline.orchestrator import run_pipeline
    from src.exporters.pptx_exporter import export_pptx
    from pptx import Presentation

    path = _make_dossier_with_2_entities(tmp_path)
    dossier = run_pipeline(inputs=[str(path)], use_llm=False, project_name="t")

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    prs = Presentation(str(out))
    dre_slide = prs.slides[6]

    tables = [s for s in dre_slide.shapes if s.has_table]
    assert len(tables) == 1, "expected exactly one summary table on DRE slide"

    tbl = tables[0].table
    # First column lists entities — header row plus 2 rows per entity
    # (Rec. Líq. + EBITDA), so for 2 entities we expect 1 + 4 = 5 rows
    assert len(tbl.rows) == 5, (
        f"expected 5 rows (1 header + 2 entities × 2 metrics); got {len(tbl.rows)}"
    )
    # Entity names should appear in column 0 of rows 1 and 3
    col0_texts = [tbl.rows[r].cells[0].text for r in range(len(tbl.rows))]
    assert "Alpha" in col0_texts
    assert "Beta" in col0_texts


def test_dre_slide_falls_back_gracefully_when_no_revenue_line(tmp_path):
    """If we somehow produce a dossier whose entities have no revenue
    line at all, the slide should render a fallback message instead
    of silently producing a slide with neither chart nor data."""
    from src.exporters.pptx_exporter import export_pptx
    from src.models.dossier import Dossier, DossierMetadata
    from src.models.company import CompanyChapter
    from src.models.financials import (
        FinancialChapter, FinancialEntity, FinancialStatement, FinancialLine,
    )
    from src.models.market import MarketChapter, TransactionChapter
    from src.models.evidence import Evidence

    # Entity with only an "Outros" line — no revenue, no EBITDA
    junk_line = FinancialLine(
        label="Outros", values={"2024": 100, "2025": 110, "2026": 120},
        is_projected={"2024": False, "2025": True, "2026": True},
        evidence=Evidence(source_file="test"),
    )
    entity = FinancialEntity(
        name="GhostCo",
        dre=FinancialStatement(lines=[junk_line], years=["2024", "2025", "2026"]),
    )
    fin = FinancialChapter(entities=[entity])
    dossier = Dossier(
        metadata=DossierMetadata(project_name="t", target_company="X"),
        company=CompanyChapter(),
        financials=fin,
        market=MarketChapter(),
        transaction=TransactionChapter(),
    )

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    from pptx import Presentation
    prs = Presentation(str(out))
    dre_slide = prs.slides[6]
    # The fallback message must be present somewhere on the slide
    all_text = " ".join(
        s.text_frame.text for s in dre_slide.shapes if s.has_text_frame
    )
    assert "Receita Líquida não identificada" in all_text


# ── P6 — Balance slide skip ─────────────────────────────────────────
def test_balance_slide_skipped_when_no_balance_data(tmp_path):
    """Dossier with DREs but no balance lines should produce a deck
    without the 'Balanço patrimonial' slide. (Regenera ships only DREs.)
    """
    from src.pipeline.orchestrator import run_pipeline
    from src.exporters.pptx_exporter import export_pptx
    from pptx import Presentation

    path = _make_dossier_with_2_entities(tmp_path)
    dossier = run_pipeline(inputs=[str(path)], use_llm=False, project_name="t")
    # The XLSX parser doesn't extract balance sheets from DRE-only
    # workbooks, so this dossier has no balance lines.
    assert not any(
        e.balance_sheet and e.balance_sheet.lines
        for e in dossier.financials.entities
    ), "test fixture should have no balance data"

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    prs = Presentation(str(out))
    titles = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    titles.append(t.split("\n")[0])
                    break  # only the title text box per slide
    joined = " | ".join(titles)
    assert "Balanço patrimonial" not in joined, (
        f"Balance slide should be skipped when no data; got titles: {joined}"
    )


def test_footer_page_numbers_renumber_on_skipped_slide(tmp_path):
    """When the balance slide is skipped, page numbers in footers
    should reflect the actual slide position (no gaps).
    """
    from src.pipeline.orchestrator import run_pipeline
    from src.exporters.pptx_exporter import export_pptx
    from pptx import Presentation

    path = _make_dossier_with_2_entities(tmp_path)
    dossier = run_pipeline(inputs=[str(path)], use_llm=False, project_name="t")

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    prs = Presentation(str(out))

    # Walk slides; for each, find the small numeric textbox near the
    # bottom-right (the page number). It must equal the slide's index.
    for idx, slide in enumerate(prs.slides, 1):
        page_num_in_footer = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t.isdigit() and len(t) <= 2:
                    page_num_in_footer = int(t)
        # Slide 1 (cover) deliberately has no footer
        if idx == 1:
            continue
        assert page_num_in_footer == idx, (
            f"Slide {idx} footer says {page_num_in_footer}, expected {idx}"
        )


# ── P7 — Overview cells ────────────────────────────────────────────
def test_overview_sheet_uses_em_dash_for_missing_values(tmp_path):
    """Empty profile fields must render as '—', not as blank cells."""
    from src.exporters.xlsx_exporter import export_xlsx
    from src.models.dossier import Dossier, DossierMetadata
    from src.models.company import CompanyChapter
    from src.models.financials import FinancialChapter
    from src.models.market import MarketChapter, TransactionChapter
    import openpyxl

    company = CompanyChapter()
    # Don't fill founding_year, headquarters, etc. — leave .value=None

    dossier = Dossier(
        metadata=DossierMetadata(project_name="P", target_company="C"),
        company=company,
        financials=FinancialChapter(),
        market=MarketChapter(),
        transaction=TransactionChapter(),
    )

    out = tmp_path / "out.xlsx"
    export_xlsx(dossier, str(out))

    wb = openpyxl.load_workbook(out, data_only=True)
    ws = wb["Visão Geral"]
    # Walk rows 3-12 (the field rows) — none should have a None at C2
    blanks = []
    for r in range(3, 13):
        v = ws.cell(row=r, column=2).value
        if v is None or v == "":
            blanks.append(r)
    assert not blanks, (
        f"rows {blanks} have blank values; should have '—' instead"
    )


def test_overview_executive_with_no_pct_renders_dash(tmp_path):
    """Executive without ownership_pct must show '—' in column C."""
    from src.exporters.xlsx_exporter import export_xlsx
    from src.models.dossier import Dossier, DossierMetadata
    from src.models.company import CompanyChapter, Executive
    from src.models.financials import FinancialChapter
    from src.models.market import MarketChapter, TransactionChapter
    import openpyxl

    company = CompanyChapter()
    company.executives = [
        Executive(name="Jane Doe", role="CEO", ownership_pct=None),
        Executive(name="John Smith", role="COO", ownership_pct=0.25),
    ]

    dossier = Dossier(
        metadata=DossierMetadata(project_name="P", target_company="C"),
        company=company,
        financials=FinancialChapter(),
        market=MarketChapter(),
        transaction=TransactionChapter(),
    )

    out = tmp_path / "out.xlsx"
    export_xlsx(dossier, str(out))

    wb = openpyxl.load_workbook(out, data_only=True)
    ws = wb["Visão Geral"]
    # Find the DIRETORIA section + the data rows
    name_cells = {
        ws.cell(row=r, column=1).value: r
        for r in range(1, ws.max_row + 1)
    }
    jane_row = name_cells["Jane Doe"]
    john_row = name_cells["John Smith"]
    # Jane has no pct → should be "—"
    assert ws.cell(row=jane_row, column=3).value == "—"
    # John has 0.25 → should be the numeric value, formatted later by Excel
    assert ws.cell(row=john_row, column=3).value == 0.25


# ── P5 — Market/Landscape placeholders ─────────────────────────────
def test_market_slide_renders_placeholder_when_no_drivers(tmp_path):
    """When no growth_drivers and no fragmentation are populated, the
    Mercado slide must still show 'Drivers de crescimento' header and a
    '—' placeholder. Otherwise the slide's bottom half goes blank.
    """
    from src.pipeline.orchestrator import run_pipeline
    from src.exporters.pptx_exporter import export_pptx
    from pptx import Presentation

    path = _make_dossier_with_2_entities(tmp_path)
    dossier = run_pipeline(inputs=[str(path)], use_llm=False, project_name="t")

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    prs = Presentation(str(out))
    # Slide 9 is Mercado (after Balance was skipped)
    market_slide = prs.slides[8]
    text_blob = " ".join(
        s.text_frame.text for s in market_slide.shapes if s.has_text_frame
    )
    # Section headers always present
    assert "Fragmentação" in text_blob
    assert "Drivers de crescimento" in text_blob
    # Placeholder for missing data
    assert ("estrutura de mercado não identificada" in text_blob
            or "drivers de crescimento não identificados" in text_blob)


def test_landscape_slide_renders_placeholder_when_no_competitors(tmp_path):
    """No competitors → explicit placeholder, not silent omission."""
    from src.pipeline.orchestrator import run_pipeline
    from src.exporters.pptx_exporter import export_pptx
    from pptx import Presentation

    path = _make_dossier_with_2_entities(tmp_path)
    dossier = run_pipeline(inputs=[str(path)], use_llm=False, project_name="t")

    out = tmp_path / "out.pptx"
    export_pptx(dossier, str(out))

    prs = Presentation(str(out))
    landscape_slide = prs.slides[9]   # after Balance skipped
    text_blob = " ".join(
        s.text_frame.text for s in landscape_slide.shapes if s.has_text_frame
    )
    assert "Concorrentes não identificados" in text_blob
    assert "Transações precedentes não identificadas" in text_blob
